import io
import traceback

def extract_text(uploaded_file) -> tuple[str, str]:
    """
    Extract plain text from any supported file type.
    Returns: (extracted_text, error_message)
    error_message is empty string if successful.
    """
    filename = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()

    if not file_bytes:
        return "", "File is empty — please upload a valid document."

    if filename.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    elif filename.endswith(".docx"):
        return _extract_docx(file_bytes)
    elif filename.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    elif filename.endswith((".txt", ".md")):
        try:
            return file_bytes.decode("utf-8", errors="ignore"), ""
        except Exception as e:
            return "", f"TXT read error: {e}"
    elif filename.endswith(".doc"):
        return "", "Legacy .doc format not supported. Please save as .docx and re-upload."
    else:
        return "", f"Unsupported file type: {filename.split('.')[-1].upper()}. Use PDF, DOCX, PPTX, or TXT."


def _extract_pdf(file_bytes: bytes) -> tuple[str, str]:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        if len(reader.pages) == 0:
            return "", "PDF has no pages."
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        text = text.strip()
        if not text:
            return "", "PDF appears to be scanned/image-based — no text could be extracted. Please use a text-based PDF."
        return text, ""
    except Exception as e:
        return "", f"PDF extraction error: {traceback.format_exc()}"


def _extract_docx(file_bytes: bytes) -> tuple[str, str]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
        text = "\n".join(parts).strip()
        if not text:
            return "", "DOCX file appears empty — no readable text found."
        return text, ""
    except Exception as e:
        return "", f"DOCX extraction error: {traceback.format_exc()}"


def _extract_pptx(file_bytes: bytes) -> tuple[str, str]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        text = "\n".join(parts).strip()
        if not text:
            return "", "PPTX file appears empty — no readable text found in slides."
        return text, ""
    except Exception as e:
        return "", f"PPTX extraction error: {traceback.format_exc()}"
